# -*- coding: utf-8 -*-
"""レクチャー・商談兼用のパワポスライドを生成するスクリプト

図解つき使い方ガイド（build_illustrated_guide.py）と同じ画面図を
スライドに流用する。生成: docs/積立金入力アシスタント_説明スライド.pptx
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
sys.path.insert(0, os.path.join(HERE, "animations"))

from build_illustrated_guide import (  # noqa: E402
    fig_overview, fig_settings, fig_menu, fig_bank_paste, fig_income,
    fig_expense, fig_roster, fig_plan, fig_reports,
)

from pptx import Presentation  # noqa: E402
from pptx.util import Cm, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

NAVY = RGBColor(0x1B, 0x2F, 0x56)
BLUE = RGBColor(0x2F, 0x5F, 0xA8)
AMBER = RGBColor(0xF5, 0xB7, 0x31)
GRAY = RGBColor(0x66, 0x70, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "游ゴシック"

# 16:9
SLIDE_W, SLIDE_H = Cm(33.867), Cm(19.05)


def _set_font(tf, size, bold=False, color=NAVY, align=None):
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def bar(slide, title_text, sub=None):
    """上部のタイトルバー"""
    box = slide.shapes.add_textbox(Cm(1.2), Cm(0.7), Cm(31.4), Cm(1.6))
    tf = box.text_frame
    tf.text = title_text
    _set_font(tf, 28, bold=True, color=NAVY)
    line = slide.shapes.add_shape(1, Cm(1.3), Cm(2.5), Cm(9.0), Cm(0.12))  # rectangle
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()
    if sub:
        sb = slide.shapes.add_textbox(Cm(1.2), Cm(2.7), Cm(31.4), Cm(1.0))
        sb.text_frame.text = sub
        _set_font(sb.text_frame, 14, color=GRAY)


def bullets(slide, items, left=Cm(1.6), top=Cm(4.2), width=Cm(30.5), size=18):
    box = slide.shapes.add_textbox(left, top, width, Cm(12))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "・" + item
        p.space_after = Pt(10)
    _set_font(tf, size)
    return box


def picture_slide(prs, title_text, img_path, note=None):
    s = add_slide(prs)
    bar(s, title_text)
    # 画像 1000x640 → 幅24cmで中央配置
    w = Cm(24)
    h = Cm(24 * 640 / 1000)
    s.shapes.add_picture(img_path, (SLIDE_W - w) / 2, Cm(3.2), width=w, height=h)
    if note:
        nb = s.shapes.add_textbox(Cm(1.6), Cm(3.3) + h + Cm(0.2), Cm(30.5), Cm(1.2))
        nb.text_frame.text = note
        _set_font(nb.text_frame, 14, color=BLUE, align=PP_ALIGN.CENTER)
    return s


def build():
    figs = {
        "overview": fig_overview(), "settings": fig_settings(), "menu": fig_menu(),
        "bank": fig_bank_paste(), "income": fig_income(), "expense": fig_expense(),
        "roster": fig_roster(), "plan": fig_plan(), "reports": fig_reports(),
    }

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- 1. タイトル ----
    s = add_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2), Cm(6.2), Cm(30), Cm(3))
    tb.text_frame.text = "積立金会計 入力アシスタント"
    _set_font(tb.text_frame, 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Cm(2), Cm(9.6), Cm(30), Cm(2))
    sb.text_frame.text = "320人分の入力を、1クリックに。"
    _set_font(sb.text_frame, 26, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    nb = s.shapes.add_textbox(Cm(2), Cm(12.4), Cm(30), Cm(1.5))
    nb.text_frame.text = "現場レクチャー ／ ご説明資料"
    _set_font(nb.text_frame, 16, color=WHITE, align=PP_ALIGN.CENTER)

    # ---- 2. 今日の内容 ----
    s = add_slide(prs)
    bar(s, "本日の内容（約60分）")
    bullets(s, [
        "これは何か ― マスターは1ミリも変えない「入力の代行係」（10分）",
        "毎月の山場を実演 ― 銀行の振替結果を貼って、ボタン1つ（15分）",
        "ハンズオン ― 皆さまの手で同じ操作を（15分）",
        "年に1回の業務 ― 320名のクラス替えが2クリック（10分）",
        "点検・保管・質疑（10分）",
    ], size=20)

    # ---- 3〜10. 図解スライド ----
    picture_slide(prs, "全体像 ― 届いたデータが、ボタン操作だけでマスターと帳票になる",
                  figs["overview"])
    picture_slide(prs, "準備は1回だけ ― 設定シートの黄色いセルを埋める",
                  figs["settings"], "マスター(C3)・年度(C5)・口座マスター(C7) の3ヶ所")
    picture_slide(prs, "業務とボタンが1対1 ― 覚えるのは12個のボタンだけ",
                  figs["menu"])
    picture_slide(prs, "毎月の山場① ― 銀行の振替結果を貼るだけで未納者が見つかる",
                  figs["bank"], "読取321件／振替済315名／未納5名 が数秒で判定される")
    picture_slide(prs, "毎月の山場② ― 未納者表は自動。⑤で全員分を一括記録",
                  figs["income"])
    picture_slide(prs, "支出があったとき ― 例外だけ書けば、記録と承認書が同時に終わる",
                  figs["expense"])
    picture_slide(prs, "毎年4月 ― 320名のクラス替えが「貼って2クリック」",
                  figs["roster"], "人が判断するのは赤い行（転入生・同姓同名）だけ")
    picture_slide(prs, "年間予定表 ― 計画は1回、実行は2クリック",
                  figs["plan"])
    picture_slide(prs, "締めの点検と帳票 ― 保管まで毎月同じ流れ",
                  figs["reports"])

    # ---- 11. 安全装置 ----
    s = add_slide(prs)
    bar(s, "「壊さない」ための4つの安全装置")
    items = [
        ("自動バックアップ", "書き込みのたびに実行前のコピーを日時付きで自動保存。いつでも戻せる"),
        ("構造チェック", "マスターの形が想定と違うファイルには書き込まない（誤爆防止）"),
        ("上書き確認", "金額が入っている列には、人数を示して確認してから書き込む"),
        ("触らない設計", "精算番号・氏名・数式列・行や列の並びは一切変更しない"),
    ]
    for i, (t, d) in enumerate(items):
        x = Cm(1.6) + (i % 2) * Cm(15.6)
        y = Cm(4.2) + (i // 2) * Cm(5.6)
        card = s.shapes.add_shape(5, x, y, Cm(15.0), Cm(5.0))  # rounded rect
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFD)
        card.line.color.rgb = BLUE
        tf = card.text_frame
        tf.word_wrap = True
        tf.text = t
        p2 = tf.add_paragraph()
        p2.text = d
        for j, p in enumerate(tf.paragraphs):
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(20 if j == 0 else 14)
                r.font.bold = (j == 0)
                r.font.color.rgb = NAVY if j == 0 else GRAY

    # ---- 12. 料金 ----
    s = add_slide(prs)
    bar(s, "料金プラン（初期費用0円・30日間無料トライアル）")
    plans = [
        ("ライト", "9,800円/月", "1学年（1会計）\nメールサポート", False),
        ("スタンダード", "19,800円/月", "全学年（3会計）\n新年度対応アップデート\n導入支援・優先サポート", True),
        ("プレミアム", "34,800円/月", "全学年＋様式カスタマイズ\n研修年2回・電話サポート\nデータ移行支援", False),
    ]
    for i, (name, price, desc, featured) in enumerate(plans):
        x = Cm(1.6) + i * Cm(10.6)
        card = s.shapes.add_shape(5, x, Cm(4.0), Cm(9.8), Cm(11.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE if not featured else RGBColor(0xFF, 0xF7, 0xE0)
        card.line.color.rgb = AMBER if featured else RGBColor(0xD5, 0xDE, 0xED)
        card.line.width = Pt(3 if featured else 1.5)
        tf = card.text_frame
        tf.word_wrap = True
        tf.text = name + ("（おすすめ）" if featured else "")
        pp = tf.add_paragraph(); pp.text = price
        for line in desc.split("\n"):
            pd = tf.add_paragraph(); pd.text = "・" + line
        for j, p in enumerate(tf.paragraphs):
            p.alignment = PP_ALIGN.CENTER if j <= 1 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = FONT
                r.font.bold = j <= 1
                r.font.size = Pt(20 if j == 0 else (28 if j == 1 else 14))
                r.font.color.rgb = NAVY if j != 1 else BLUE
    nb = s.shapes.add_textbox(Cm(1.6), Cm(16.0), Cm(30.5), Cm(1.6))
    nb.text_frame.text = "年間契約・税別／年払いは2ヶ月分割引／請求書払い対応。生徒1人あたり月約62円（320名校・スタンダード）"
    _set_font(nb.text_frame, 14, color=GRAY, align=PP_ALIGN.CENTER)

    # ---- 13. クロージング ----
    s = add_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2), Cm(6.5), Cm(30), Cm(3))
    tb.text_frame.text = "今年のクラス替えから、間に合います。"
    _set_font(tb.text_frame, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Cm(2), Cm(10.2), Cm(30), Cm(2))
    sb.text_frame.text = "初期費用0円・30日間無料トライアル｜まずは架空データの練習環境からノーリスクで"
    _set_font(sb.text_frame, 18, color=AMBER, align=PP_ALIGN.CENTER)
    cb = s.shapes.add_textbox(Cm(2), Cm(13.0), Cm(30), Cm(1.5))
    cb.text_frame.text = "お問い合わせ: info@cocorolab.co.jp"
    _set_font(cb.text_frame, 16, color=WHITE, align=PP_ALIGN.CENTER)

    out = os.path.join(HERE, "積立金入力アシスタント_説明スライド.pptx")
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    build()
