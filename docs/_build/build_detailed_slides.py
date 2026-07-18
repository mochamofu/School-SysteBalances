# -*- coding: utf-8 -*-
"""詳細解説スライド（60分版・約42枚）を生成する。

14枚の商談用スライド（build_slides.py）とは別に、機能①〜⑮を
1機能1枚で解説する「研修・詳細説明用」のロングバージョン。
各スライドに発表者ノート（話すことのメモ）を埋め込む。

生成: docs/02_営業・商談資料/積立金入力アシスタント_詳細解説スライド60分.pptx
"""
import os
import sys

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)

import build_slides as bs  # noqa: E402  ヘルパーと図を流用
from build_illustrated_guide import (  # noqa: E402
    fig_overview, fig_settings, fig_menu, fig_bank_paste, fig_income,
    fig_expense, fig_roster, fig_plan, fig_reports,
)
from pptx import Presentation  # noqa: E402
from pptx.util import Cm, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

NAVY, BLUE, AMBER, GRAY, WHITE = bs.NAVY, bs.BLUE, bs.AMBER, bs.GRAY, bs.WHITE
FONT = bs.FONT
SLIDE_W, SLIDE_H = bs.SLIDE_W, bs.SLIDE_H
GREEN = RGBColor(0x1E, 0x7A, 0x3C)

OUT = os.path.join(HERE, "..", "02_営業・商談資料", "積立金入力アシスタント_詳細解説スライド60分.pptx")


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def divider(prs, part, title, minutes):
    s = bs.add_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2), Cm(6.4), Cm(30), Cm(2))
    tb.text_frame.text = part
    bs._set_font(tb.text_frame, 22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    tb2 = s.shapes.add_textbox(Cm(2), Cm(8.4), Cm(30), Cm(3))
    tb2.text_frame.text = title
    bs._set_font(tb2.text_frame, 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb3 = s.shapes.add_textbox(Cm(2), Cm(12.2), Cm(30), Cm(1.4))
    tb3.text_frame.text = f"（約{minutes}分）"
    bs._set_font(tb3.text_frame, 18, color=WHITE, align=PP_ALIGN.CENTER)
    return s


def stat_slide(prs, title, stats, note_text=None):
    """大きな数字3つのスライド。stats = [(数字, ラベル), ...]"""
    s = bs.add_slide(prs)
    bs.bar(s, title)
    n = len(stats)
    card_w = Cm(9.4)
    gap = (SLIDE_W - card_w * n) / (n + 1)
    for i, (num, label) in enumerate(stats):
        x = gap * (i + 1) + card_w * i
        card = s.shapes.add_shape(5, x, Cm(5.0), card_w, Cm(8.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFD)
        card.line.color.rgb = BLUE; card.line.width = Pt(1.5)
        tf = card.text_frame; tf.word_wrap = True
        tf.text = num
        p2 = tf.add_paragraph(); p2.text = label
        for j, p in enumerate(tf.paragraphs):
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = FONT
                r.font.bold = (j == 0)
                r.font.size = Pt(48 if j == 0 else 15)
                r.font.color.rgb = BLUE if j == 0 else NAVY
    if note_text:
        nb = s.shapes.add_textbox(Cm(1.6), Cm(14.0), Cm(30.5), Cm(1.4))
        nb.text_frame.text = note_text
        bs._set_font(nb.text_frame, 15, color=GRAY, align=PP_ALIGN.CENTER)
    return s


def feature_slide(prs, num_label, title, when, ops, result, caution, img=None):
    """機能詳解スライド: 左に4項目、右に図（あれば）"""
    s = bs.add_slide(prs)
    bs.bar(s, f"{num_label}　{title}")
    text_w = Cm(17.5) if img else Cm(30.5)
    box = s.shapes.add_textbox(Cm(1.6), Cm(3.4), text_w, Cm(14.5))
    tf = box.text_frame; tf.word_wrap = True
    rows = [("いつ使う", when), ("操作", ops), ("起きること", result), ("注意・コツ", caution)]
    first = True
    for head, body in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "■ " + head
        p.space_before = Pt(6)
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BLUE
        for line in body:
            p2 = tf.add_paragraph()
            p2.text = "　" + line
            p2.space_after = Pt(2)
            for r in p2.runs:
                r.font.name = FONT; r.font.size = Pt(14.5); r.font.color.rgb = NAVY
    if img:
        w = Cm(13.5)
        h = Cm(13.5 * 640 / 1000)
        s.shapes.add_picture(img, Cm(19.6), Cm(4.6), width=w, height=h)
    return s


def build():
    figs = {
        "overview": fig_overview(), "settings": fig_settings(), "menu": fig_menu(),
        "bank": fig_bank_paste(), "income": fig_income(), "expense": fig_expense(),
        "roster": fig_roster(), "plan": fig_plan(), "reports": fig_reports(),
    }
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ---- 1 タイトル ----
    s = bs.add_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2), Cm(6.0), Cm(30), Cm(3))
    tb.text_frame.text = "積立金会計 入力アシスタント　詳細解説"
    bs._set_font(tb.text_frame, 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Cm(2), Cm(9.6), Cm(30), Cm(2))
    sb.text_frame.text = "しくみ・全機能・安全装置・導入の進め方（60分版）"
    bs._set_font(sb.text_frame, 22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    notes(s, "挨拶。今日は詳しい版。前半でしくみ、真ん中で15個の機能を全部、後半で安全装置と導入の進め方。"
             "質問はいつでも遮ってOKと伝える。使うデータはすべて架空と宣言。")

    # ---- 2 本日の流れ ----
    s = bs.add_slide(prs)
    bs.bar(s, "本日の流れ（約60分）")
    bs.bullets(s, [
        "第1部　いまの業務と課題の整理（5分）",
        "第2部　製品の全体像 ― マスターに触らないしくみ（10分）",
        "第3部　初期設定 ― 最初の1回だけやること（5分）",
        "第4部　15個の機能を全部見る（25分）",
        "第5部　安全装置と年間の運用（10分）",
        "第6部　導入の進め方・質疑（5分）",
    ], size=20)
    notes(s, "時間配分を先に見せて安心してもらう。第4部が本体。途中休憩は入れない代わりに質問随時。")

    # ================= 第1部 =================
    notes(divider(prs, "第1部", "いまの業務と課題の整理", 5),
          "まず現状の確認から。参加者の実務と一致するか、頷きを確認しながら進める。")

    s = bs.add_slide(prs)
    bs.bar(s, "いまの毎月の流れ（確認）")
    bs.bullets(s, [
        "ゆうちょ銀行から振替結果が届く（紙 または データ）",
        "振替できなかった生徒を目で探す",
        "全員分の入金額を、1人ずつマスターへ手入力",
        "支出が発生するたび、同じ金額を全員分入力し、例外の生徒を1人ずつ修正",
        "同じ内容を支出承認書・収入承認書にもう一度書く",
        "年度末は電卓で項目ごとに集計して決算書へ",
    ], size=19)
    notes(s, "「この流れで合っていますか？」と必ず聞く。違う点があればメモして後の説明で拾う。")

    stat_s = stat_slide(prs, "数字で見る、いまの負担", [
        ("320人分", "毎回の入力対象（全校）"),
        ("月5〜10時間", "入力・照合・転記の手作業"),
        ("1件", "誤請求は1件でも保護者対応に発展"),
    ], "時間の負担と、金額の大小によらない「間違えられない」緊張の両方が課題")
    notes(stat_s, "時間だけでなく心理的負担を強調。預かり金だから1円のミスも説明責任がある、という共感を取る。")

    s = bs.add_slide(prs)
    bs.bar(s, "いまのExcelマスターを「変えない」理由")
    bs.bullets(s, [
        "長年使われてきたマスターは、様式も数式も現場と監査に馴染んだ資産",
        "新システムへの移行は、データ載せ替え・二重運用・研修のコストが大きい",
        "だから本製品は、マスターを置き換えない。1ミリも変えない",
        "変えるのは「入力のしかた」だけ ― 手入力を、ボタンに任せる",
    ], size=20)
    notes(s, "既存のやり方を否定しないことが大事。マスターは正しい、入力だけが辛い、という整理。")

    # ================= 第2部 =================
    notes(divider(prs, "第2部", "製品の全体像", 10),
          "ここから製品の話。キーワードは『隣に置く』『読むのと書くのを分ける』。")

    notes(bs.picture_slide(prs, "全体像 ― 届いたデータが、ボタン操作だけでマスターと帳票になる",
                           figs["overview"]),
          "矢印を左から右へ。届くデータ→アシスタント→マスター・帳票。アシスタントが「入力の代行係」。")

    s = bs.add_slide(prs)
    bs.bar(s, "お渡しするファイルはこれだけ")
    bs.bullets(s, [
        "積立金入力アシスタント.xlsm ― 本体。ボタン①〜⑮が入ったExcelブック",
        "口座マスターひな形.xlsx ― 精算番号と、ゆうちょ口座の対応表（⑪で使用）",
        "練習用データ一式 ― 架空の80名。何度でも安全に練習できる",
        "検証用データ一式 ― 架空の320名。実物と同じ規模での確認用",
        "サーバー・インストール・アカウント・ネット接続 ― すべて不要",
    ], size=19)
    notes(s, "「ソフトを入れる」のではなく「Excelファイルを置く」。ここが公立校で大事なポイント。"
             "詳しくは管理者向け説明資料があると案内。")

    s = bs.add_slide(prs)
    bs.bar(s, "「マスターに触らない」6つの約束")
    bs.bullets(s, [
        "行・列の挿入・削除・並べ替えは一切しない（精算書の数式を壊さない）",
        "精算番号・氏名の列には書き込まない",
        "数式の列（未納印・合計・残金など）には書き込まない",
        "書き込むのは、金額のセルと見出しだけ",
        "書き込む前に、必ずバックアップ（控えのコピー）を自動作成",
        "マスターの形が想定と違うファイルには、書き込み自体を拒否",
    ], size=19)
    notes(s, "6つ読み上げる価値あり。監査・引き継ぎの観点で刺さる。「壊しようがない設計」と言い切る。")

    s = bs.add_slide(prs)
    bs.bar(s, "パソコンに対して「行わないこと」")
    bs.bullets(s, [
        "インストーラーの実行 ― なし（フォルダにコピーするだけ）",
        "管理者権限 ― 不要",
        "レジストリ・システム設定の変更 ― なし",
        "常駐プログラム ― なし（Excelを閉じれば何も動かない）",
        "ネットワーク通信 ― なし（完全オフライン・外部送信ゼロ）",
        "削除するとき ― ファイルを消すだけ",
    ], size=19)
    notes(s, "情報担当が同席していたらこのスライドを厚めに。管理者向けA4を配布資料として渡せる。")

    # ================= 第3部 =================
    notes(divider(prs, "第3部", "初期設定 ― 最初の1回だけ", 5),
          "設定は3ヶ所だけ、と先に言ってから見せる。")

    notes(bs.picture_slide(prs, "設定シート ― 黄色いセルを3ヶ所埋めるだけ",
                           figs["settings"], "マスターの場所（C3）・年度（C5）・口座マスターの場所（C7）"),
          "導入日は一緒に入れるので、覚える必要はないと安心させる。")

    s = bs.add_slide(prs)
    bs.bar(s, "口座マスターを1回だけ作る")
    bs.bullets(s, [
        "精算番号・氏名・口座記号・口座番号 の4列を、ひな形に入力（400行まで）",
        "⑪振替結果の照合が、この表を使って「どの口座＝どの生徒」を判定する",
        "全角の数字が混ざっていても自動でそろえて照合する",
        "兄弟で同じ口座を共有している場合は「要確認」として人間に返す（勝手に判定しない）",
        "作るのは1回だけ。翌年は転出入の差分を直すだけ",
    ], size=19)
    notes(s, "共有口座の扱いは質問が出やすい。機械が勝手に決めず人に返す設計、と答える。")

    notes(bs.picture_slide(prs, "練習環境 ― 架空の80名でいつでも練習できる",
                           figs["menu"], "実在の生徒情報はどこにも入っていません"),
          "壊しても作り直せる練習環境がある、が導入の心理的ハードルを下げる。本日のデモも全部これ。")

    # ================= 第4部: 機能詳解 =================
    notes(divider(prs, "第4部", "15個の機能を全部見る", 25),
          "1機能1枚×15。各スライド90秒目安。デモを挟むなら①④⑪の3ヶ所。")

    F = feature_slide
    notes(F(prs, "①", "名簿を解析して照合する",
            ["毎年4月、クラス替えの掲示用名簿が出たとき"],
            ["掲示用名簿をシートにそのまま貼り付けて、ボタンを押すだけ"],
            ["クラスのブロック配置を自動で見つけ、氏名でマスターと照合",
             "全員の新しい組・番号が一覧になる（320名なら320行）"],
            ["転入生・同姓同名だけ赤い行になる ― 人が見るのはそこだけ",
             "少人数クラス（3名以上）も自動で検出"],
            figs["roster"]),
          "「何日かかっていましたか？」と聞いてから。赤い行しか見なくていい、が核心。")

    notes(F(prs, "②", "クラス替えをマスターに反映する",
            ["①の照合結果を確認した後"],
            ["ボタン1つ。書き込み前に自動バックアップ"],
            ["全員分の組・番号がマスターに一括反映される"],
            ["精算番号と氏名は変更しない（並び順が崩れない）",
             "反映後も①の一覧が残るので、後から見直せる"],
            figs["roster"]),
          "①と②で「見る」と「書く」を分けている＝確認してから書く設計、と説明。")

    notes(F(prs, "③", "新入生としてマスターに登録する",
            ["新年度、1年生のマスターを新しく作るとき"],
            ["名簿を貼って①で解析 → ③で空のマスターに登録"],
            ["精算番号1番から順に、組・番号・氏名が入った新しいマスターができる"],
            ["空のマスター（様式だけのファイル）はこちらで用意して納品",
             "3年間このマスターを使い続ける（学年進行）"],
            figs["roster"]),
          "1年生用。入学時に1回だけ。以後3年間は②のクラス替えだけで回る。")

    notes(F(prs, "④", "支出をマスターへ一括入力",
            ["業者への支払いが発生するたび（遠足・教材・検定など）"],
            ["支出No・件名・日付・一人あたり金額を入力",
             "対象外や金額が違う生徒だけ「例外表」に精算番号で書く → ボタン"],
            ["在籍全員に金額が入り、例外だけ上書きされる",
             "同じ内容で支出承認書シートも自動で埋まる（印刷するだけ）"],
            ["転退学者は0、給付型は0、返金はマイナスで書く",
             "存在しない精算番号を書くとエラーで止まる（黙って無視しない）"],
            figs["expense"]),
          "一番使うボタン。例外だけ書く、という発想の転換を丁寧に。315人入力→例外3人だけ書く。")

    notes(F(prs, "⑤", "収入をマスターへ一括入力",
            ["口座振替の結果が確定したとき（毎月）"],
            ["収入枠No・件名・一人あたり金額を入力",
             "振替できなかった生徒だけ「未納者表」に精算番号 → ボタン"],
            ["未納者以外の全員に入金額が入る",
             "未納者は空欄のまま → マスターのH列に未納の印が自動で立つ",
             "収入承認書も自動で埋まる"],
            ["未納者表は⑪が自動で作るので、通常は手入力不要"],
            figs["income"]),
          "⑪とセットで使う。未納の印が数式で自動、という既存マスターの仕組みを活かしている点を強調。")

    notes(F(prs, "⑥", "収入枠の一覧を表示",
            ["⑤で使う枠Noを決めるとき"],
            ["ボタン1つ"],
            ["43個の収入枠の使用状況（項目名・人数）が一覧で出る"],
            ["空いている枠がすぐ分かる。上書きミスの予防"],
            figs["income"]),
          "小さい機能だが「どの枠を使えばいいか」の迷いを消す。30秒でよい。")

    notes(F(prs, "⑦", "決算用の集計を実行",
            ["年度末、決算書を作るとき"],
            ["ボタン1つ"],
            ["全項目の 人数・一人あたり・執行総額 が一覧になる",
             "この数字を決算書へ転記するだけ"],
            ["電卓での項目別集計が不要になる",
             "「一人あたり」は最も多くの生徒に入っている金額を自動判定"],
            figs["reports"]),
          "年度末の残業がここで消える。実データ検証で合計が1円も違わないことを確認済みと言える。")

    notes(F(prs, "⑧", "マスターの整合性をチェック",
            ["月次の締め・決算前・不安なとき、いつでも"],
            ["ボタン1つ"],
            ["収支の突き合わせ・未納者一覧・入力漏れの点検結果が出る"],
            ["書き込みは一切しない読み取り専用の点検",
             "月1回の実行を習慣にすると監査対応が楽"],
            figs["reports"]),
          "「健康診断ボタン」。押しても何も変わらないから毎月押してよい、と伝える。")

    notes(F(prs, "⑨", "精算書を一括PDF保存",
            ["精算書を配る・保管するとき"],
            ["ボタン1つ"],
            ["生徒ごとのPDF（組-番号_氏名.pdf）が全員分できる"],
            ["紙の一括印刷しかなかった精算書を、デジタルでも保存できる",
             "保存先は設定シートで変更可能"],
            figs["reports"]),
          "320枚のPDFが数分でできる。保護者説明・保管・再発行が楽になる。")

    notes(F(prs, "⑩", "精算書を一括印刷",
            ["精算書を紙で配るとき"],
            ["ボタン1つ（範囲指定も可能）"],
            ["全員分の精算書が順に印刷される"],
            ["マスターに元々入っていた1999年製の印刷マクロを、安全に作り直したもの",
             "動きは今までと同じなので、使い勝手が変わらない"],
            figs["reports"]),
          "既存マクロの後継と説明すると安心される。「作った人がいなくても、もう直せる」の実例。")

    notes(F(prs, "⑪", "振替結果を照合",
            ["ゆうちょから振替結果が届いたとき（毎月の山場）"],
            ["結果の 記号・番号・金額・結果 を貼り付けてボタン"],
            ["口座マスターと照合し、生徒ごとに 振替済／未納 を自動判定",
             "未納者は⑤の未納者表へ自動転記",
             "読取件数・済・未納・不明の集計も表示"],
            ["不明口座・共有口座・重複行は「要確認」として人間に返す",
             "1000行まで一度に貼れる（全校分OK）"],
            figs["bank"]),
          "本製品の心臓。目視で探す作業がゼロになる。デモをやるならここ。数字を予告してから実行。")

    notes(F(prs, "⑫", "予定を入力フォームへ転送",
            ["年間予定表に書いた予定を実行するとき"],
            ["予定の行番号を指定するだけ"],
            ["支出入力／収入入力のフォームに件名・金額が自動で入る",
             "あとは④または⑤を押すだけ"],
            ["毎回請求書から件名や金額を打ち直さなくてよくなる"],
            figs["plan"]),
          "計画は年度初めに1回、実行は2クリック、というリズムを見せる。")

    notes(F(prs, "⑬", "支出項目を読み込む",
            ["年度末・年度替わりに、1年間の支出を振り返るとき"],
            ["ボタン1つ（マスターは読むだけ・書き換えない）"],
            ["支出100枠が一覧になる（件名・日付・人数・一人あたり・総額）",
             "業者名の列は自由に書き換えられる（次の⑭で使う）"],
            ["未使用の枠は自動で除外される",
             "書き込んだ業者名やメモは、読み直しても消えない"],
            figs["plan"]),
          "⑬⑭⑮は年度替わりの3点セット。ここからの3枚は続き物として話す。")

    notes(F(prs, "⑭", "業者別に集計する",
            ["「この業者に年間いくら使ったか」を知りたいとき"],
            ["⑬の一覧で業者名をそろえて、ボタン"],
            ["業者ごとの 項目数・年間合計・月別内訳（4月〜3月）が出る",
             "同じ業者が複数の列に分かれていても1行に合算される"],
            ["例：教材業者が5列に分散 → 業者名をそろえるだけで年間合計が出る"],
            figs["plan"]),
          "実データ検証で、5列に分かれた業者の合算が1円も違わないことを確認済み。")

    notes(F(prs, "⑮", "来年度の予定へ引き継ぐ",
            ["翌年度のファイルを準備するとき"],
            ["⑬の一覧で、来年も使う項目に○を付けてボタン"],
            ["○の項目だけが年間予定表に自動で追加される"],
            ["「今年これ要らない」→ ×にするだけ",
             "「今年これ追加したい」→ 予定表に1行書くだけ",
             "毎年のファイル作りが○×の付け替えで済む"],
            figs["plan"]),
          "項目は毎年入れ替わる、という現場の実情に対する答え。3年間の学年進行がこのサイクルで回る。")

    # ================= 第5部 =================
    notes(divider(prs, "第5部", "安全装置と年間の運用", 10),
          "機能の次は「壊さない仕組み」。事務室で一番大事な話。")

    s = bs.add_slide(prs)
    bs.bar(s, "4つの安全装置（すべて自動で働く）")
    items = [
        ("自動バックアップ", "書き込みのたびに、実行前のコピーを日時付きで自動保存。いつでも戻せる"),
        ("構造チェック", "マスターの形が想定と違うファイルには書き込まない（誤ったファイルへの誤爆防止）"),
        ("上書き確認", "すでに金額が入っている列には、人数を示して確認してから書き込む"),
        ("入力チェック", "存在しない精算番号・空欄の金額は、その場でエラーにして止まる"),
    ]
    for i, (t, dsc) in enumerate(items):
        x = Cm(1.6) + (i % 2) * Cm(15.6)
        y = Cm(4.2) + (i // 2) * Cm(5.6)
        card = s.shapes.add_shape(5, x, y, Cm(15.0), Cm(5.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFD)
        card.line.color.rgb = BLUE
        tf = card.text_frame; tf.word_wrap = True
        tf.text = t
        p2 = tf.add_paragraph(); p2.text = dsc
        for j, p in enumerate(tf.paragraphs):
            for r in p.runs:
                r.font.name = FONT
                r.font.size = Pt(20 if j == 0 else 14)
                r.font.bold = (j == 0)
                r.font.color.rgb = NAVY if j == 0 else GRAY
    notes(s, "4枚のカードを1枚ずつ。バックアップは「毎回・自動・同じPC内」の3語で。")

    s = bs.add_slide(prs)
    bs.bar(s, "わざと間違えると、どうなるか")
    bs.bullets(s, [
        "存在しない精算番号（例：400）を書いて実行 → 「範囲外です」で停止。何も書かれない",
        "金額欄を空のまま実行 → 「金額欄が空です」で停止",
        "設定のファイル場所が空のまま実行 → 入力を促す案内で停止",
        "マスター以外のExcelを指定して実行 → 「形が想定と違います」で書き込み拒否",
        "壊れる前に、止まる ― これを異常データ15項目のテストで確認済み",
    ], size=19)
    notes(s, "「間違えたらどうなるの」への先回り回答。検証ログがあることも一言添える。")

    s = bs.add_slide(prs)
    bs.bar(s, "年間カレンダー ― いつ、どのボタンを使うか")
    bs.bullets(s, [
        "4月 ― 名簿貼付→①照合→②反映（新1年は③登録）。年間予定表を記入",
        "毎月 ― 振替結果を貼る→⑪照合→⑤一括入力。支出のたび④",
        "毎月の締め ― ⑧整合性チェック → マスターと帳票を保管用PCへ",
        "年度末 ― ⑦決算集計 → 決算書へ転記。⑨精算書PDF／⑩印刷",
        "年度替わり ― ⑬読み込み→業者名と○×→⑭業者別集計→⑮翌年度へ引き継ぎ",
    ], size=19)
    notes(s, "1年分を1枚で。これが配布資料にもなっている（使い方ガイド参照）。")

    notes(bs.picture_slide(prs, "毎月の締め ― 点検と帳票の保管まで同じ流れ",
                           figs["reports"], "マスター＋精算書PDF＋チェック結果をワンセットで保管用PCへ"),
          "月次ルーチンの完成形。属人化せず、誰がやっても同じ流れになる。")

    # ================= 第6部 =================
    notes(divider(prs, "第6部", "導入の進め方・質疑", 5),
          "最後に導入の話。ここは押し売りせず、段階を見せて安心してもらう。")

    s = bs.add_slide(prs)
    bs.bar(s, "導入は3段階 ― いきなり本番には使いません")
    bs.bullets(s, [
        "第1段階　練習（架空の80名）― 壊してよい環境でひととおり操作に慣れる",
        "第2段階　並走（実マスターのコピー）― 手入力の結果と突き合わせて、1円も違わないことを自分の目で確認",
        "第3段階　本番切り替え ― 設定シートの場所を本物に変えるだけ。以後も毎回自動バックアップ",
        "各段階に動作確認チェックシート（期待値つき）を用意",
        "うまくいかないときは、いつでも元のやり方に戻れる（マスターは無傷のまま）",
    ], size=19)
    notes(s, "「いきなり本番に使わない」を強調。並走期間が信頼をつくる。チェックシートを見せる。")

    s = bs.add_slide(prs)
    bs.bar(s, "まとめ ― 今日の3点")
    bs.bullets(s, [
        "いまのマスターは1ミリも変えない。変わるのは入力のしかただけ",
        "毎月の山場（振替照合）と年1回の山場（クラス替え）が、貼ってボタンになる",
        "壊さない安全装置と、練習→並走→本番の段階導入で、安心して始められる",
    ], size=22)
    notes(s, "3点だけ復唱。ここで質問を受ける。")

    # ---- クロージング ----
    s = bs.add_slide(prs)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = s.shapes.add_textbox(Cm(2), Cm(6.5), Cm(30), Cm(3))
    tb.text_frame.text = "ご質問をどうぞ"
    bs._set_font(tb.text_frame, 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sb = s.shapes.add_textbox(Cm(2), Cm(10.2), Cm(30), Cm(2))
    sb.text_frame.text = "料金は学校のご予算に合わせて個別にご相談｜まずは無料の練習環境から"
    bs._set_font(sb.text_frame, 18, color=AMBER, align=PP_ALIGN.CENTER)
    notes(s, "料金はここでも深追いしない。トライアル開始日の相談へつなげる。")

    prs.save(os.path.abspath(OUT))
    print(f"saved {os.path.abspath(OUT)}  ({len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)}枚)")


if __name__ == "__main__":
    build()
